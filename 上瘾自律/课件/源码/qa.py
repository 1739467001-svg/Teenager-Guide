#!/usr/bin/env python3
# 几何 QA：越界 / 文字溢出 / 文本框重叠 / 边距
# 中文字形宽度接近 1em，估算比拉丁文可靠得多
import sys, math, unicodedata
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
SLIDE_W, SLIDE_H = 13.3, 7.5
MARGIN_MIN = 0.45          # 允许 kicker/页码略靠边

def cw(ch):
    """字符宽度，单位 em"""
    if ch == '\n': return 0
    e = unicodedata.east_asian_width(ch)
    if e in ('W', 'F'): return 1.0          # 全角中文
    if ch == ' ': return 0.30
    return 0.52                              # 半角

def para_text(p):
    return ''.join(r.text for r in p.runs)

def est_lines(text, box_w_in, pt):
    """估算这段文字要占几行"""
    em = pt / 72.0
    if em <= 0 or box_w_in <= 0: return 0
    cap = box_w_in / em                      # 一行能放多少 em
    total = 0
    for seg in text.split('\n'):
        w = sum(cw(c) for c in seg)
        total += max(1, math.ceil(w / cap)) if w > 0 else 1
    return total

def analyze(path):
    prs = Presentation(path)
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            if sh.left is None or sh.top is None: continue
            x, y = sh.left / EMU_IN, sh.top / EMU_IN
            w, h = (sh.width or 0) / EMU_IN, (sh.height or 0) / EMU_IN
            name = (sh.shape_type or '?')

            # 1) 越界
            if x < -0.02 or y < -0.02 or x + w > SLIDE_W + 0.02 or y + h > SLIDE_H + 0.02:
                # 背景图铺满整页，跳过
                if not (abs(x) < 0.05 and abs(y) < 0.05 and w >= SLIDE_W - 0.1):
                    issues.append((si, 'OUT', f'越界 x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}'))

            if not sh.has_text_frame: continue
            tf = sh.text_frame
            txt = '\n'.join(para_text(p) for p in tf.paragraphs).strip()
            if not txt: continue

            # 字号：取该框最大 run
            pts = [r.font.size.pt for p in tf.paragraphs for r in p.runs if r.font.size]
            pt = max(pts) if pts else 18.0
            # 行距
            ls = None
            for p in tf.paragraphs:
                if p.line_spacing:
                    ls = p.line_spacing if isinstance(p.line_spacing, float) else p.line_spacing.pt
                    break
            lh = (ls if (ls and ls > 4) else pt * 1.42) / 72.0

            lines = est_lines(txt, w, pt)
            need = lines * lh
            if need > h * 1.06 and h > 0:
                issues.append((si, 'OVERFLOW',
                    f'"{txt[:26]}…" {pt:.0f}pt 需 {need:.2f}" / 框高 {h:.2f}" ({lines}行)'))

            # 2) 边距
            if 0 <= x < MARGIN_MIN or (x + w) > SLIDE_W - MARGIN_MIN + 0.02:
                issues.append((si, 'MARGIN', f'"{txt[:20]}…" 距边 {min(x, SLIDE_W-(x+w)):.2f}"'))

            boxes.append((x, y, w, h, txt[:20], pt))

        # 3) 文本框重叠
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                ox = min(a[0]+a[2], b[0]+b[2]) - max(a[0], b[0])
                oy = min(a[1]+a[3], b[1]+b[3]) - max(a[1], b[1])
                if ox > 0.12 and oy > 0.12:
                    issues.append((si, 'OVERLAP', f'"{a[4]}…" ∩ "{b[4]}…" 重叠 {ox:.2f}×{oy:.2f}"'))
    return len(prs.slides), issues

for path in sys.argv[1:]:
    n, iss = analyze(path)
    print(f"\n{'='*72}\n{path.split('/')[-1]}   {n} 页")
    if not iss:
        print("  ✓ 无越界 / 无溢出 / 无重叠 / 边距合格")
    else:
        from collections import Counter
        print(f"  发现 {len(iss)} 项：{dict(Counter(k for _,k,_ in iss))}")
        for si, kind, msg in iss:
            print(f"   p{si:>3} [{kind:<8}] {msg}")
